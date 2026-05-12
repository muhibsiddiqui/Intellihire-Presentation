"""
IntelliHire FYP Presentation Generator
Dark navy theme — mirrors the website colour palette exactly.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Palette (matches website) ───────────────────────────────────────────────
NAVY_950  = RGBColor(0x05, 0x0c, 0x1a)   # slide background
NAVY_800  = RGBColor(0x0d, 0x1b, 0x33)   # card background
NAVY_600  = RGBColor(0x1a, 0x2d, 0x4a)   # card border / secondary bg
BRAND_400 = RGBColor(0x2f, 0x97, 0xf7)   # primary blue
CYAN_400  = RGBColor(0x00, 0xd4, 0xff)   # secondary cyan
EMERALD   = RGBColor(0x10, 0xb9, 0x81)   # green
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)   # amber
PURPLE    = RGBColor(0xa8, 0x55, 0xf7)   # purple
RED       = RGBColor(0xef, 0x44, 0x44)   # red
WHITE     = RGBColor(0xff, 0xff, 0xff)
WHITE_70  = RGBColor(0xb0, 0xc4, 0xde)
WHITE_50  = RGBColor(0x8a, 0xa0, 0xb8)

# ─── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ─── Helpers ─────────────────────────────────────────────────────────────────
def blank_slide():
    layout = prs.slide_layouts[6]          # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color=NAVY_950):
    """Fill the slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(1)):
    """Add a rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def label(slide, text, l, t, w, h,
          size=Pt(11), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox

def heading_badge(slide, badge_text, title_text, subtitle_text="",
                  title_color=WHITE):
    """Common top-of-slide header block."""
    # Badge pill
    b = box(slide, Inches(0.5), Inches(0.35), Inches(2.6), Inches(0.32),
            fill=RGBColor(0x1a, 0x3a, 0x5c), line_color=BRAND_400, line_width=Pt(0.75))
    label(slide, badge_text, Inches(0.52), Inches(0.35), Inches(2.56), Inches(0.32),
          size=Pt(8), bold=True, color=BRAND_400, align=PP_ALIGN.CENTER)

    # Title
    label(slide, title_text, Inches(0.5), Inches(0.72), Inches(12.33), Inches(0.65),
          size=Pt(30), bold=True, color=title_color, align=PP_ALIGN.LEFT)

    # Subtitle
    if subtitle_text:
        label(slide, subtitle_text, Inches(0.5), Inches(1.35), Inches(12.33), Inches(0.38),
              size=Pt(12), color=WHITE_70, align=PP_ALIGN.LEFT)

def accent_bar(slide):
    """Thin horizontal accent line under the header."""
    bar = box(slide, Inches(0.5), Inches(1.78), Inches(12.33), Inches(0.025),
              fill=BRAND_400)

def card(slide, l, t, w, h, accent_color=BRAND_400):
    """Glassmorphism-style card."""
    c = box(slide, l, t, w, h, fill=NAVY_800, line_color=accent_color, line_width=Pt(0.75))
    return c

def dot_bullet(slide, items, l, t, w, col=EMERALD, size=Pt(10.5), spacing_h=Inches(0.26)):
    """Render a list of strings as dot-prefixed bullet lines."""
    y = t
    for item in items:
        dot = box(slide, l, y + Inches(0.07), Inches(0.1), Inches(0.1), fill=col)
        label(slide, item, l + Inches(0.16), y, w - Inches(0.18), spacing_h,
              size=size, color=WHITE_70)
        y += spacing_h
    return y

def stat_box(slide, l, t, w, h, value, title, color=BRAND_400):
    card(slide, l, t, w, h, accent_color=color)
    label(slide, value, l, t + Inches(0.12), w, Inches(0.38),
          size=Pt(24), bold=True, color=color, align=PP_ALIGN.CENTER)
    label(slide, title, l, t + Inches(0.5), w, Inches(0.28),
          size=Pt(9), color=WHITE_70, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)

# Subtle grid-like dots via very small boxes (decorative)
for xi in range(0, 14):
    for yi in range(0, 8):
        dot = box(s, Inches(xi), Inches(yi), Inches(0.04), Inches(0.04),
                  fill=RGBColor(0x1a, 0x2a, 0x40))

# Glow circle (decorative)
glow = box(s, Inches(3), Inches(1.5), Inches(7), Inches(5),
           fill=RGBColor(0x06, 0x14, 0x26), line_color=None)

# FYP badge
box(s, Inches(4.4), Inches(0.55), Inches(4.5), Inches(0.38),
    fill=RGBColor(0x0d, 0x2a, 0x48), line_color=BRAND_400, line_width=Pt(0.75))
label(s, "● FINAL YEAR PROJECT — 2025-2026", Inches(4.4), Inches(0.54), Inches(4.5), Inches(0.4),
      size=Pt(9), bold=True, color=BRAND_400, align=PP_ALIGN.CENTER)

# INTELLIHIRE title
label(s, "INTELLI", Inches(2.5), Inches(1.1), Inches(4.2), Inches(1.2),
      size=Pt(72), bold=True, color=BRAND_400, align=PP_ALIGN.RIGHT)
label(s, "HIRE", Inches(6.65), Inches(1.1), Inches(4.2), Inches(1.2),
      size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
label(s, "AI-Powered Candidate Screening & Interview Automation System",
      Inches(1.5), Inches(2.45), Inches(10.33), Inches(0.5),
      size=Pt(16), color=WHITE_70, align=PP_ALIGN.CENTER)

# Divider
box(s, Inches(4.5), Inches(3.05), Inches(4.33), Inches(0.02), fill=BRAND_400)

# Live URL callout
box(s, Inches(3.8), Inches(3.25), Inches(5.7), Inches(0.42),
    fill=RGBColor(0x0a, 0x22, 0x3a), line_color=CYAN_400, line_width=Pt(1))
label(s, "🌐  Live at  intellihire.com.pk", Inches(3.8), Inches(3.24), Inches(5.7), Inches(0.44),
      size=Pt(13), bold=True, color=CYAN_400, align=PP_ALIGN.CENTER)

# Team left
label(s, "PROJECT TEAM", Inches(1.5), Inches(3.85), Inches(4), Inches(0.28),
      size=Pt(8), bold=True, color=RGBColor(0x4a, 0x7a, 0xa0), align=PP_ALIGN.LEFT)
team = [
    "22K-4418   Muhammad Omer Khan",
    "22K-4428   Muhib Siddiqui",
    "22K-4403   Syed Muhammad Salik Ahmed",
]
y = Inches(4.15)
for t_ in team:
    label(s, t_, Inches(1.5), y, Inches(4.5), Inches(0.28), size=Pt(11), color=WHITE_70)
    y += Inches(0.28)

# Project details right
label(s, "PROJECT DETAILS", Inches(7.3), Inches(3.85), Inches(4), Inches(0.28),
      size=Pt(8), bold=True, color=RGBColor(0x4a, 0x7a, 0xa0), align=PP_ALIGN.LEFT)
details = [
    "Supervisor:   Miss Sobia Iftikhar",
    "Type:            Final Year Project (2025-2026)",
    "Institution:   FAST NUCES Karachi",
]
y = Inches(4.15)
for d in details:
    label(s, d, Inches(7.3), y, Inches(5), Inches(0.28), size=Pt(11), color=WHITE_70)
    y += Inches(0.28)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "The Challenge", "Problem Statement",
              "Identifying the core inefficiencies IntelliHire was designed to solve")
accent_bar(s)

problems = [
    ("Recruitment Bottlenecks", BRAND_400, [
        "Manual screening is slow, inconsistent, and bias-prone",
        "HR teams spend 70% of time on initial screening tasks",
        "Candidates wait weeks before a single interaction",
        "Limited evaluation depth from manual shortlisting",
    ]),
    ("HR Information Access", AMBER, [
        "Employees depend on HR staff for basic policy queries",
        "No centralized, searchable HR document knowledge base",
        "No secure authenticated access for personalized retrieval",
    ]),
    ("The System Gap", RED, [
        "No integrated platform combines AI screening with an HR assistant",
        "No scalable system for automated interviews with live proctoring",
        "Companies use expensive enterprise tools or fully manual processes",
    ]),
]

col_w = Inches(3.9)
col_gap = Inches(0.17)
col_top = Inches(2.0)
col_h = Inches(4.95)

for i, (title, color, pts) in enumerate(problems):
    lx = Inches(0.5) + i * (col_w + col_gap)
    card(s, lx, col_top, col_w, col_h, accent_color=color)
    # Colored top bar
    box(s, lx, col_top, col_w, Inches(0.06), fill=color)
    label(s, title, lx + Inches(0.15), col_top + Inches(0.12), col_w - Inches(0.3), Inches(0.38),
          size=Pt(13), bold=True, color=color)
    dot_bullet(s, pts, lx + Inches(0.15), col_top + Inches(0.58), col_w - Inches(0.25),
               col=color, size=Pt(10.5), spacing_h=Inches(0.29))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OBJECTIVES
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "Project Goals", "Key Objectives",
              "Core goals that define the scope and purpose of IntelliHire")
accent_bar(s)

objectives = [
    ("AI-Driven Recruitment Automation", BRAND_400, [
        "End-to-end AI screening from CV upload to final evaluation report",
        "Voice-based AI interviews with real-time dynamic follow-up questions",
        "Multi-modal analysis — voice, face, gaze — for holistic assessment",
    ]),
    ("Intelligent HR Assistance", CYAN_400, [
        "RAG-based HR chatbot for instant policy and document queries",
        "Secure employee self-service for personalized knowledge retrieval",
        "Reduce HR workload by automating routine information requests",
    ]),
    ("Comprehensive Analytics & Integrity", EMERALD, [
        "Detailed candidate reports with multi-dimensional scoring metrics",
        "Data-driven feedback to streamline HR decision-making",
        "Real-time proctoring ensuring fair and tamper-free sessions",
    ]),
]

col_w = Inches(3.9)
col_gap = Inches(0.17)
col_top = Inches(2.0)
col_h = Inches(4.95)

for i, (title, color, pts) in enumerate(objectives):
    lx = Inches(0.5) + i * (col_w + col_gap)
    card(s, lx, col_top, col_w, col_h, accent_color=color)
    box(s, lx, col_top, col_w, Inches(0.06), fill=color)
    label(s, title, lx + Inches(0.15), col_top + Inches(0.12), col_w - Inches(0.3), Inches(0.45),
          size=Pt(13), bold=True, color=color)
    dot_bullet(s, pts, lx + Inches(0.15), col_top + Inches(0.65), col_w - Inches(0.25),
               col=color, size=Pt(10.5), spacing_h=Inches(0.34))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "System Design", "Architecture Overview",
              "Four-layer architecture powering intelligent recruitment automation")
accent_bar(s)

layers = [
    ("L1  PRESENTATION LAYER", "React 19 · TypeScript · Material UI · Socket.IO · WebRTC",
     ["Interview UI", "HR Dashboard", "Candidate Portal", "Analytics"], BRAND_400),
    ("L2  APPLICATION LAYER", "Flask 3.0 · JWT Auth · SQLAlchemy · Flask-CORS · Gunicorn",
     ["Interview Engine", "CV / HR Services", "Proctoring Pipeline"], CYAN_400),
    ("L3  AI / ML LAYER", "DeepSeek · LangChain · ChromaDB · YOLOv8 · MediaPipe",
     ["STT / TTS", "RAG + Embeddings", "Gaze & Object Detection"], PURPLE),
    ("L4  DATA LAYER", "MySQL (Relational) · ChromaDB (Vector Store)",
     ["User & Interview Data", "Embeddings & Documents"], EMERALD),
]

connectors = ["REST + WebSocket", "LangChain / API Calls", "ORM + Vector Queries"]
layer_h = Inches(0.9)
layer_w = Inches(12.33)
top_start = Inches(1.95)
gap = Inches(0.16)

for i, (title, tech, modules, color) in enumerate(layers):
    lt = top_start + i * (layer_h + gap)
    card(s, Inches(0.5), lt, layer_w, layer_h, accent_color=color)
    box(s, Inches(0.5), lt, Inches(0.08), layer_h, fill=color)

    # Layer number badge
    box(s, Inches(0.65), lt + Inches(0.22), Inches(0.5), Inches(0.45),
        fill=RGBColor(0x0d, 0x1b, 0x33), line_color=color, line_width=Pt(0.75))
    label(s, title[:2], Inches(0.65), lt + Inches(0.22), Inches(0.5), Inches(0.45),
          size=Pt(10), bold=True, color=color, align=PP_ALIGN.CENTER)

    label(s, title[4:], Inches(1.25), lt + Inches(0.12), Inches(3.5), Inches(0.36),
          size=Pt(12), bold=True, color=color)
    label(s, tech, Inches(1.25), lt + Inches(0.48), Inches(5.2), Inches(0.32),
          size=Pt(9), color=WHITE_50, italic=True)

    # Modules as chips
    mx = Inches(6.7)
    for mod in modules:
        mw = Inches(1.7) if len(mod) < 14 else Inches(2.4)
        box(s, mx, lt + Inches(0.28), mw, Inches(0.34),
            fill=RGBColor(0x0a, 0x18, 0x2e), line_color=color, line_width=Pt(0.5))
        label(s, mod, mx, lt + Inches(0.28), mw, Inches(0.34),
              size=Pt(8.5), color=color, align=PP_ALIGN.CENTER)
        mx += mw + Inches(0.12)

    # Connector arrow between layers
    if i < len(connectors):
        arrow_t = lt + layer_h + Inches(0.015)
        box(s, Inches(6.1), arrow_t, Inches(1.3), Inches(0.14), fill=RGBColor(0x08, 0x14, 0x28))
        label(s, f"↕  {connectors[i]}", Inches(5.5), arrow_t - Inches(0.01), Inches(2.5), Inches(0.16),
              size=Pt(7.5), color=WHITE_50, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "Technical Implementation", "Technology Stack",
              "Production-grade technologies powering every layer of the platform")
accent_bar(s)

backend = [
    ("Python / Flask 3.0",         "REST API + WebSocket server",                   BRAND_400),
    ("DeepSeek API",                "LLM engine — Q&A generation, scoring",          RGBColor(0x4a, 0x90, 0xd9)),
    ("LangChain",                   "AI orchestration & RAG pipeline",               EMERALD),
    ("ChromaDB",                    "Vector store for CV & JD embeddings",           RED),
    ("YOLOv8",                      "Object detection for proctoring",               AMBER),
    ("MySQL",                       "Primary relational DB (SQLAlchemy ORM)",        RGBColor(0x44, 0x79, 0xa1)),
    ("MediaPipe + OpenCV",          "Gaze tracking & facial landmark detection",     CYAN_400),
    ("SpeechRecognition + gTTS",    "STT / TTS pipeline for voice interviews",       EMERALD),
]

frontend = [
    ("React 19",                    "Modern UI library with hooks & context",        CYAN_400),
    ("TypeScript",                  "Type-safe frontend development",                RGBColor(0x31, 0x78, 0xc6)),
    ("Material UI (MUI)",           "Component library & theming system",            BRAND_400),
    ("Socket.IO Client",            "Real-time WebSocket communication",             EMERALD),
    ("WebRTC",                      "Browser-based video & audio capture",           RGBColor(0xf0, 0x62, 0x92)),
    ("Recharts",                    "Data visualization & analytics charts",         RED),
    ("Axios",                       "HTTP client for REST API integration",          PURPLE),
    ("React Router v6",             "Client-side routing & navigation",              RED),
]

for col_i, (group_title, techs, header_color) in enumerate([
        ("Backend Technologies", backend, BRAND_400),
        ("Frontend Technologies", frontend, CYAN_400)]):
    lx = Inches(0.5) + col_i * Inches(6.5)
    label(s, group_title, lx, Inches(1.95), Inches(6.0), Inches(0.32),
          size=Pt(12), bold=True, color=header_color)
    box(s, lx, Inches(2.27), Inches(6.0), Inches(0.025), fill=header_color)

    row_h = Inches(0.53)
    for j, (name, desc, color) in enumerate(techs):
        ty = Inches(2.35) + j * row_h
        card(s, lx, ty, Inches(6.1), row_h - Inches(0.06), accent_color=color)
        box(s, lx, ty, Inches(0.06), row_h - Inches(0.06), fill=color)
        label(s, name, lx + Inches(0.15), ty + Inches(0.05), Inches(2.5), Inches(0.24),
              size=Pt(10), bold=True, color=WHITE)
        label(s, desc, lx + Inches(0.15), ty + Inches(0.26), Inches(5.7), Inches(0.22),
              size=Pt(8.5), color=WHITE_50)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — INTERVIEW FLOW
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "Core Feature — Module 1", "AI Interview Flow",
              "Six-stage pipeline from interviewer setup to scored final report")
accent_bar(s)

steps = [
    ("01", "Interviewer Setup",         "HR defines job description, scoring weights, evaluation criteria, must-ask questions & behavioral constraints.",          BRAND_400),
    ("02", "CV Upload & RAG",           "Candidate uploads CV → text extraction → Sentence Transformer embeddings → ChromaDB → semantic JD matching.",             CYAN_400),
    ("03", "Question Generation",       "DeepSeek generates tailored questions from CV-JD match, interviewer params, role requirements via RAG pipeline.",          RGBColor(0x58, 0xac, 0xee)),
    ("04", "Live AI Interview",         "Voice interview via STT/TTS over WebSocket. Sub-200ms audio segmentation. Multi-turn dialogue with sliding context.",      BRAND_400),
    ("05", "Real-time Monitoring",      "YOLOv8 detects people & phones. MediaPipe tracks gaze. Browser Visibility API logs tab switches — all timestamped.",       AMBER),
    ("06", "Scoring & Report",          "Multi-modal scoring across 5 dimensions via DeepSeek. Full report with sentiment timeline, proctoring log & feedback.",    EMERALD),
]

cols = 3
rows = 2
cell_w = Inches(4.1)
cell_h = Inches(2.2)
col_gap = Inches(0.17)
row_gap = Inches(0.12)
start_x = Inches(0.5)
start_y = Inches(2.0)

for idx, (num, title, desc, color) in enumerate(steps):
    col = idx % cols
    row = idx // cols
    lx = start_x + col * (cell_w + col_gap)
    ly = start_y + row * (cell_h + row_gap)

    card(s, lx, ly, cell_w, cell_h, accent_color=color)
    box(s, lx, ly, cell_w, Inches(0.05), fill=color)

    # Step number
    box(s, lx + Inches(0.15), ly + Inches(0.12), Inches(0.5), Inches(0.5),
        fill=RGBColor(0x05, 0x10, 0x22), line_color=color, line_width=Pt(1))
    label(s, num, lx + Inches(0.15), ly + Inches(0.12), Inches(0.5), Inches(0.5),
          size=Pt(13), bold=True, color=color, align=PP_ALIGN.CENTER)

    label(s, title, lx + Inches(0.75), ly + Inches(0.15), cell_w - Inches(0.9), Inches(0.42),
          size=Pt(12), bold=True, color=color)
    label(s, desc, lx + Inches(0.15), ly + Inches(0.68), cell_w - Inches(0.3), Inches(1.38),
          size=Pt(9.5), color=WHITE_70)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — REAL-TIME PROCTORING
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "Anti-Cheating System", "Real-Time Proctoring",
              "Multi-layered integrity monitoring for fair and secure interview sessions")
accent_bar(s)

proctoring = [
    ("Gaze Tracking",          "MediaPipe eye landmark detection flags sustained off-screen gaze at 30 fps with configurable thresholds.",        "30 fps tracking",      BRAND_400),
    ("Mobile Detection",       "YOLOv8 detects phones, tablets, and secondary screens with bounding-box overlay on every video frame.",           "YOLOv8 inference",     AMBER),
    ("Tab Monitoring",         "Browser Visibility API records every focus-loss event with timestamp & duration. Candidates are immediately flagged.", "Event-based",      RED),
    ("Multi-Person Detection", "YOLOv8 person detection identifies if multiple individuals are present. Immediate flag with frame capture for evidence.", "Real-time",    PURPLE),
]

cell_w = Inches(6.05)
cell_h = Inches(2.3)
col_gap = Inches(0.15)
row_gap = Inches(0.12)
start_x = Inches(0.5)
start_y = Inches(2.0)

for idx, (title, desc, stat, color) in enumerate(proctoring):
    col = idx % 2
    row = idx // 2
    lx = start_x + col * (cell_w + col_gap)
    ly = start_y + row * (cell_h + row_gap)

    card(s, lx, ly, cell_w, cell_h, accent_color=color)
    box(s, lx, ly, cell_w, Inches(0.05), fill=color)

    # Stat chip
    chip_w = Inches(1.3)
    box(s, lx + cell_w - chip_w - Inches(0.15), ly + Inches(0.12), chip_w, Inches(0.3),
        fill=RGBColor(0x05, 0x10, 0x22), line_color=color, line_width=Pt(0.5))
    label(s, stat, lx + cell_w - chip_w - Inches(0.15), ly + Inches(0.12), chip_w, Inches(0.3),
          size=Pt(8), color=color, align=PP_ALIGN.CENTER)

    label(s, title, lx + Inches(0.2), ly + Inches(0.12), cell_w - Inches(1.8), Inches(0.38),
          size=Pt(14), bold=True, color=color)
    label(s, desc, lx + Inches(0.2), ly + Inches(0.58), cell_w - Inches(0.4), Inches(1.6),
          size=Pt(10.5), color=WHITE_70)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MULTI-MODAL SCORING
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "Evaluation Engine", "Multi-Modal Scoring",
              "Five-dimensional candidate evaluation powered by DeepSeek AI")
accent_bar(s)

dimensions = [
    ("Verbal Confidence",   "Lexical hedging, filler-word density, assertiveness markers in speech",                                     "20%",  BRAND_400),
    ("Semantic Relevance",  "Cosine similarity + DeepSeek judgment score against rubric expected answers",                               "30%",  CYAN_400),
    ("Sentiment Analysis",  "Real-time sentiment via speech transcription and facial micro-expression timeline",                         "15%",  EMERALD),
    ("Proctoring Score",    "Gaze tracking, mobile/person detection, tab monitoring aggregated anomaly score",                           "15%",  AMBER),
    ("CV Relevance",        "ATS scoring: skills match, experience, education, projects, certifications",                                "20%",  PURPLE),
]

dim_w = Inches(2.38)
dim_h = Inches(4.5)
dim_gap = Inches(0.14)
dim_top = Inches(2.05)

for i, (name, desc, weight, color) in enumerate(dimensions):
    lx = Inches(0.5) + i * (dim_w + dim_gap)
    card(s, lx, dim_top, dim_w, dim_h, accent_color=color)
    box(s, lx, dim_top, dim_w, Inches(0.06), fill=color)

    # Weight badge
    box(s, lx + dim_w - Inches(0.65), dim_top + Inches(0.1), Inches(0.55), Inches(0.38),
        fill=RGBColor(0x05, 0x10, 0x22), line_color=color, line_width=Pt(0.75))
    label(s, weight, lx + dim_w - Inches(0.65), dim_top + Inches(0.1), Inches(0.55), Inches(0.38),
          size=Pt(12), bold=True, color=color, align=PP_ALIGN.CENTER)

    label(s, name, lx + Inches(0.12), dim_top + Inches(0.58), dim_w - Inches(0.24), Inches(0.55),
          size=Pt(11), bold=True, color=color)
    label(s, desc, lx + Inches(0.12), dim_top + Inches(1.22), dim_w - Inches(0.24), Inches(3.1),
          size=Pt(9.5), color=WHITE_70)

# Weight bar
bar_top = dim_top + dim_h + Inches(0.18)
bar_x = Inches(0.5)
total = Inches(12.33)
for i, (_, _, weight, color) in enumerate(dimensions):
    w_frac = int(weight.strip('%')) / 100
    seg_w = int(total * w_frac)
    box(s, bar_x, bar_top, seg_w - 4, Inches(0.22), fill=color)
    label(s, weight, bar_x, bar_top, seg_w - 4, Inches(0.22),
          size=Pt(8), bold=True, color=NAVY_950, align=PP_ALIGN.CENTER)
    bar_x += seg_w

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — HR CHATBOT
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "Module 2 — HR Assistant", "Intelligent HR Chatbot",
              "RAG-based document Q&A system — instant access to policies, procedures, and knowledge")
accent_bar(s)

chatbot_features = [
    ("📁", "Document Management",    "Upload PDF, DOCX, TXT → auto-extract → semantic chunking → vector storage",                  BRAND_400),
    ("🔍", "RAG Retrieval",           "Sentence Transformers (all-MiniLM-L6-v2) → ChromaDB similarity search → context injection",   CYAN_400),
    ("🤖", "AI Response Generation",  "Retrieved context + query → Gemini AI → cited, accurate answers",                             EMERALD),
    ("💬", "Multi-Turn Conversations","Session-based history tracking with intent classification and context awareness",              PURPLE),
    ("🔐", "Role-Based Access",       "Admin/HR can upload documents — all employees can query the knowledge base securely",         AMBER),
]

feat_w = Inches(5.8)
feat_h = Inches(0.88)
feat_gap = Inches(0.1)
feat_x = Inches(0.5)

for i, (icon, title, desc, color) in enumerate(chatbot_features):
    ty = Inches(2.0) + i * (feat_h + feat_gap)
    card(s, feat_x, ty, feat_w, feat_h, accent_color=color)
    box(s, feat_x, ty, feat_w, Inches(0.05), fill=color)
    # Icon area
    box(s, feat_x + Inches(0.12), ty + Inches(0.16), Inches(0.55), Inches(0.55),
        fill=RGBColor(0x0a, 0x18, 0x2e), line_color=color, line_width=Pt(0.5))
    label(s, icon, feat_x + Inches(0.12), ty + Inches(0.14), Inches(0.55), Inches(0.55),
          size=Pt(16), align=PP_ALIGN.CENTER)
    label(s, title, feat_x + Inches(0.82), ty + Inches(0.1), feat_w - Inches(1.0), Inches(0.32),
          size=Pt(11), bold=True, color=color)
    label(s, desc, feat_x + Inches(0.82), ty + Inches(0.42), feat_w - Inches(1.0), Inches(0.38),
          size=Pt(9.5), color=WHITE_70)

# RAG pipeline diagram on the right
rx = Inches(6.65)
ry = Inches(2.0)
rw = Inches(6.3)
rh = Inches(5.3)
card(s, rx, ry, rw, rh, accent_color=BRAND_400)
label(s, "RAG Pipeline Flow", rx + Inches(0.2), ry + Inches(0.15), rw - Inches(0.4), Inches(0.35),
      size=Pt(12), bold=True, color=BRAND_400)
box(s, rx + Inches(0.2), ry + Inches(0.5), rw - Inches(0.4), Inches(0.025), fill=BRAND_400)

rag_steps = [
    ("HR Uploads Document",             BRAND_400),
    ("Text Extraction & Chunking",      CYAN_400),
    ("Sentence Transformer Embeddings", PURPLE),
    ("ChromaDB Vector Storage",         RED),
    ("Employee Query",                  AMBER),
    ("Similarity Search + LLM Answer",  EMERALD),
]
sy = ry + Inches(0.65)
for txt, col in rag_steps:
    box(s, rx + Inches(0.3), sy, rw - Inches(0.6), Inches(0.5),
        fill=RGBColor(0x05, 0x10, 0x22), line_color=col, line_width=Pt(0.75))
    label(s, txt, rx + Inches(0.3), sy, rw - Inches(0.6), Inches(0.5),
          size=Pt(10.5), color=col, align=PP_ALIGN.CENTER)
    if txt != rag_steps[-1][0]:
        box(s, rx + rw/2 - Inches(0.05), sy + Inches(0.52), Inches(0.1), Inches(0.14), fill=col)
    sy += Inches(0.66)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CURRENT STATUS
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "Where We Stand", "Current Status",
              "Core system complete — deployed live on Contabo VPS — final QA and handover phase")
accent_bar(s)

# Stats row
stats = [
    ("40+", "Backend API Routes",   BRAND_400),
    ("60+", "React Components",     CYAN_400),
    ("5",   "AI Models Integrated", PURPLE),
    ("5",   "Scoring Dimensions",   EMERALD),
]
stat_w = Inches(2.9)
stat_h = Inches(0.9)
for i, (val, lbl, col) in enumerate(stats):
    stat_box(s, Inches(0.5) + i * Inches(3.05), Inches(2.02),
             stat_w, stat_h, val, lbl, col)

# Two columns below
completed = [
    "Full-stack application (Flask + React + MySQL)",
    "AI interview engine with STT / TTS pipeline",
    "RAG-based dynamic question generation (ChromaDB)",
    "Real-time proctoring (YOLOv8 + MediaPipe + Gaze)",
    "Multi-modal scoring & report generation (DeepSeek)",
    "CV upload, parsing & ATS scoring (5 dimensions)",
    "HR dashboard with interview & candidate management",
    "Candidate portal with slot booking & results view",
    "Marketing landing page with pricing & onboarding",
    "Role-based access control (Admin / HR / Candidate)",
    "Deployed on Contabo VPS — live at intellihire.com.pk ✓",
    "Production Monitoring — uptime, logs, resource stability on live Contabo VPS",
    "End-to-End Testing — full QA across all portals, edge cases, stress tests",
    "Demo & Handover — supervisor demo, documentation, FYP submission package",
]

remaining: list = []

lx1 = Inches(0.5)
lx2 = Inches(6.92)
top_c = Inches(3.05)
col_h = Inches(4.28)
col_w = Inches(6.1)

card(s, lx1, top_c, col_w, col_h, accent_color=EMERALD)
label(s, "● Completed Features", lx1 + Inches(0.15), top_c + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
      size=Pt(11), bold=True, color=EMERALD)
dot_bullet(s, completed, lx1 + Inches(0.2), top_c + Inches(0.48), col_w - Inches(0.35),
           col=EMERALD, size=Pt(9), spacing_h=Inches(0.3))

card(s, lx2, top_c, col_w, col_h, accent_color=AMBER)
label(s, "● Future Enhancements", lx2 + Inches(0.15), top_c + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
      size=Pt(11), bold=True, color=AMBER)
ry2 = top_c + Inches(0.5)
for title, desc in remaining:
    box(s, lx2 + Inches(0.15), ry2 + Inches(0.05), Inches(0.32), Inches(0.32),
        fill=RGBColor(0x0a, 0x18, 0x28), line_color=AMBER, line_width=Pt(0.75))
    label(s, title, lx2 + Inches(0.58), ry2, col_w - Inches(0.72), Inches(0.3),
          size=Pt(10.5), bold=True, color=AMBER)
    label(s, desc, lx2 + Inches(0.58), ry2 + Inches(0.3), col_w - Inches(0.72), Inches(0.32),
          size=Pt(9.5), color=WHITE_70)
    ry2 += Inches(0.75)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DEPLOYMENT ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "Infrastructure", "Deployment Architecture",
              "Live production infrastructure running on Contabo VPS — serving intellihire.com.pk")
accent_bar(s)

# Domain / CI CD bar
top_boxes = [
    ("Vercel",          "Landing Page + Marketing",  BRAND_400, Inches(0.5)),
    ("GitHub Actions",  "CI/CD Pipeline",             CYAN_400,  Inches(4.6)),
    ("intellihire.com.pk", "Domain + Email",           PURPLE,    Inches(8.7)),
]
for lbl, sub, col, lx in top_boxes:
    card(s, lx, Inches(2.0), Inches(4.0), Inches(0.9), accent_color=col)
    label(s, lbl, lx, Inches(2.05), Inches(4.0), Inches(0.36),
          size=Pt(13), bold=True, color=col, align=PP_ALIGN.CENTER)
    label(s, sub, lx, Inches(2.41), Inches(4.0), Inches(0.28),
          size=Pt(9), color=WHITE_50, align=PP_ALIGN.CENTER)

# Arrow
box(s, Inches(6.5), Inches(2.92), Inches(0.35), Inches(0.5), fill=CYAN_400)
label(s, "↓  SSH + git pull + build", Inches(5.0), Inches(3.05), Inches(3.3), Inches(0.28),
      size=Pt(8.5), color=WHITE_50, align=PP_ALIGN.CENTER, italic=True)

# VPS Container
vps_top = Inches(3.48)
vps_h = Inches(3.55)
card(s, Inches(0.5), vps_top, Inches(12.33), vps_h, accent_color=BRAND_400)
label(s, "Contabo Cloud VPS 20 — Singapore Region",
      Inches(0.5), vps_top + Inches(0.1), Inches(12.33), Inches(0.36),
      size=Pt(13), bold=True, color=BRAND_400, align=PP_ALIGN.CENTER)
label(s, "4 vCPU · 8 GB RAM · 200 GB NVMe · ~$6-8/mo",
      Inches(0.5), vps_top + Inches(0.46), Inches(12.33), Inches(0.28),
      size=Pt(9), color=WHITE_50, align=PP_ALIGN.CENTER, italic=True)

vps_services = [
    ("Nginx",           "Reverse Proxy",       BRAND_400),
    ("Flask + PM2",     "API Server (24/7)",   CYAN_400),
    ("MySQL",           "Database",            AMBER),
    ("React Build",     "Static Assets",       EMERALD),
    ("YOLOv8 + AI",     "ML Inference",        RED),
    ("ChromaDB",        "Vector Store",        RGBColor(0xff, 0x6b, 0x6b)),
    ("APScheduler",     "Email Automation",    PURPLE),
    ("Netdata / UptimeRobot", "Monitoring",    EMERALD),
]
sw = Inches(1.4)
sh = Inches(0.68)
sg = Inches(0.18)
row1 = vps_services[:4]
row2 = vps_services[4:]

for row, ry in [(row1, vps_top + Inches(0.88)), (row2, vps_top + Inches(1.68))]:
    total_row_w = len(row) * (sw + sg) - sg
    sx = Inches(0.5) + (Inches(12.33) - total_row_w) / 2
    for svc, sub, col in row:
        box(s, sx, ry, sw, sh, fill=NAVY_600, line_color=col, line_width=Pt(0.5))
        label(s, svc, sx, ry + Inches(0.1), sw, Inches(0.28),
              size=Pt(9.5), bold=True, color=col, align=PP_ALIGN.CENTER)
        label(s, sub, sx, ry + Inches(0.38), sw, Inches(0.22),
              size=Pt(8), color=WHITE_50, align=PP_ALIGN.CENTER)
        sx += sw + sg

# Stats bottom row
facts = [
    ("intellihire.com.pk", "Live URL", "Production deployment publicly accessible",    CYAN_400),
    ("5 min alert",        "Uptime Monitoring", "UptimeRobot pings every 5 minutes",   BRAND_400),
    ("PM2 managed",        "Auto Recovery", "Auto-restart, survives reboots",           EMERALD),
]
fw = Inches(4.0)
for i, (val, title, desc, col) in enumerate(facts):
    fx = Inches(0.5) + i * Inches(4.17)
    stat_box(s, fx, vps_top + vps_h + Inches(0.12), fw, Inches(0.88), val, title, col)
    label(s, desc, fx, vps_top + vps_h + Inches(0.82), fw, Inches(0.25),
          size=Pt(8), color=WHITE_50, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PRODUCT VISION
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
heading_badge(s, "The Journey Ahead", "Product Vision",
              "Transforming IntelliHire from a Final Year Project into a commercially viable SaaS product")
accent_bar(s)

# Journey steps
journey = [
    ("💡", "Idea",    "Problem identified"),
    ("🎓", "FYP",     "Built the prototype"),
    ("🚀", "Product", "Full SaaS architecture"),
    ("📈", "Scale",   "Deploy & grow"),
]
jw = Inches(2.8)
jgap = Inches(0.5)
jx = Inches(0.8)
jy = Inches(2.05)
for k, (icon, lbl, desc) in enumerate(journey):
    col = BRAND_400 if k < 3 else EMERALD
    card(s, jx, jy, jw, Inches(1.2), accent_color=col)
    label(s, icon, jx, jy + Inches(0.12), jw, Inches(0.5), size=Pt(22), align=PP_ALIGN.CENTER)
    label(s, lbl, jx, jy + Inches(0.62), jw, Inches(0.3), size=Pt(12), bold=True, color=col, align=PP_ALIGN.CENTER)
    label(s, desc, jx, jy + Inches(0.9), jw, Inches(0.26), size=Pt(9), color=WHITE_50, align=PP_ALIGN.CENTER)
    if k < len(journey) - 1:
        box(s, jx + jw, jy + Inches(0.55), jgap, Inches(0.06), fill=BRAND_400)
    jx += jw + jgap

# Pillars
pillars = [
    ("TARGET MARKET",   "Companies conducting high-volume candidate screening who need to scale hiring without scaling HR teams.", BRAND_400),
    ("BUSINESS MODEL",  "Subscription-based SaaS, tiered by interviews per month. Manual onboarding initially, then self-serve.", CYAN_400),
    ("GO-TO-MARKET",    "Landing page → Lead capture → Manual outreach → Onboarding → Automated emails → Self-serve dashboard.", EMERALD),
]
pw = Inches(3.9)
pgap = Inches(0.17)
px = Inches(0.5)
py = Inches(3.55)
ph = Inches(3.65)
for title, text, col in pillars:
    card(s, px, py, pw, ph, accent_color=col)
    box(s, px, py, pw, Inches(0.06), fill=col)
    label(s, title, px + Inches(0.15), py + Inches(0.14), pw - Inches(0.3), Inches(0.3),
          size=Pt(9), bold=True, color=col)
    label(s, text, px + Inches(0.15), py + Inches(0.52), pw - Inches(0.3), Inches(3.0),
          size=Pt(10.5), color=WHITE_70)
    px += pw + pgap

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)

# Central glow
for i in range(3):
    opacity_fill = RGBColor(0x05 + i*2, 0x0e + i*4, 0x1e + i*6)
    box(s, Inches(1.5) - Inches(i * 0.5), Inches(0.5) - Inches(i * 0.3),
        Inches(10.33) + Inches(i), Inches(6.5) + Inches(i * 0.6), fill=opacity_fill)

# Badge
box(s, Inches(4.9), Inches(0.42), Inches(3.5), Inches(0.36),
    fill=RGBColor(0x0d, 0x2a, 0x48), line_color=BRAND_400, line_width=Pt(0.75))
label(s, "Final Year Project", Inches(4.9), Inches(0.41), Inches(3.5), Inches(0.38),
      size=Pt(9), bold=True, color=BRAND_400, align=PP_ALIGN.CENTER)

# Title
label(s, "Project", Inches(1.0), Inches(0.9), Inches(11.33), Inches(0.7),
      size=Pt(36), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
label(s, "Complete & Live", Inches(1.0), Inches(1.58), Inches(11.33), Inches(0.7),
      size=Pt(36), bold=True, color=BRAND_400, align=PP_ALIGN.CENTER)

label(s, "IntelliHire is fully deployed and working — engineered to production standards, now running live.",
      Inches(1.5), Inches(2.28), Inches(10.33), Inches(0.36),
      size=Pt(12), color=WHITE_70, align=PP_ALIGN.CENTER)

# Key points
points = [
    "Fully functional FYP built end-to-end with real-world architecture — not a prototype or proof-of-concept",
    "Covers the complete hiring pipeline: job posting, CV screening, AI interviews, proctoring, scoring, HR dashboards",
    "Deep full-stack capability: React, Flask, AI/ML integration, WebRTC, DevOps, and live cloud deployment on Contabo VPS",
    "Now deployed and working at intellihire.com.pk with production infrastructure and domain routing",
]
py2 = Inches(2.75)
for pt in points:
    box(s, Inches(2.2), py2 + Inches(0.06), Inches(0.12), Inches(0.12),
        fill=BRAND_400)
    label(s, pt, Inches(2.45), py2, Inches(8.7), Inches(0.32),
          size=Pt(10.5), color=WHITE_70)
    py2 += Inches(0.38)

# Live Production badge
box(s, Inches(4.3), Inches(4.45), Inches(4.7), Inches(0.5),
    fill=RGBColor(0x08, 0x1f, 0x38), line_color=BRAND_400, line_width=Pt(1.5))
label(s, "● Live Production FYP", Inches(4.3), Inches(4.44), Inches(4.7), Inches(0.52),
      size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Now Live callout
card(s, Inches(3.2), Inches(5.1), Inches(6.9), Inches(1.4), accent_color=CYAN_400)
box(s, Inches(3.2), Inches(5.1), Inches(6.9), Inches(0.05), fill=CYAN_400)
label(s, "Now Live", Inches(3.35), Inches(5.16), Inches(6.6), Inches(0.3),
      size=Pt(10), bold=True, color=CYAN_400)
label(s, "IntelliHire is deployed on Contabo Cloud VPS and accessible at intellihire.com.pk.\nThe next milestone is onboarding real companies and candidates while continuing production monitoring.",
      Inches(3.35), Inches(5.46), Inches(6.6), Inches(0.95),
      size=Pt(10), color=WHITE_70)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)

for i in range(4):
    box(s, Inches(2) - Inches(i * 0.8), Inches(1) - Inches(i * 0.5),
        Inches(9.33) + Inches(i * 1.6), Inches(5.5) + Inches(i), fill=RGBColor(0x05 + i, 0x0c + i*3, 0x1a + i*5))

label(s, "Thank You!", Inches(0.5), Inches(1.3), Inches(12.33), Inches(1.3),
      size=Pt(60), bold=True, color=BRAND_400, align=PP_ALIGN.CENTER)
label(s, "We appreciate your time and welcome any questions.",
      Inches(0.5), Inches(2.7), Inches(12.33), Inches(0.42),
      size=Pt(16), color=WHITE_70, align=PP_ALIGN.CENTER)

# Live URL highlight
box(s, Inches(3.8), Inches(3.3), Inches(5.7), Inches(0.55),
    fill=RGBColor(0x08, 0x1f, 0x38), line_color=CYAN_400, line_width=Pt(1.5))
label(s, "🌐  intellihire.com.pk", Inches(3.8), Inches(3.29), Inches(5.7), Inches(0.57),
      size=Pt(18), bold=True, color=CYAN_400, align=PP_ALIGN.CENTER)

box(s, Inches(2.0), Inches(4.0), Inches(9.33), Inches(0.02), fill=RGBColor(0x1a, 0x2d, 0x4a))

label(s, "IntelliHire — FAST NUCES Karachi — FYP 2025-2026",
      Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.38),
      size=Pt(14), color=WHITE_70, align=PP_ALIGN.CENTER)
label(s, "Muhammad Omer Khan  ·  Muhib Siddiqui  ·  Syed Muhammad Salik Ahmed",
      Inches(0.5), Inches(4.62), Inches(12.33), Inches(0.36),
      size=Pt(12), color=WHITE_50, align=PP_ALIGN.CENTER)
label(s, "Supervised by Miss Sobia Iftikhar",
      Inches(0.5), Inches(4.98), Inches(12.33), Inches(0.32),
      size=Pt(11), color=WHITE_50, align=PP_ALIGN.CENTER)

# ─── SAVE ────────────────────────────────────────────────────────────────────
out = r"d:\FYP PPT WEB\IntelliHire_FYP_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
